"""Tests for PptxParser

Tests slide text, table, and speaker notes extraction from .pptx files.
Tests structure mode JSON metadata output with shape IDs and positions.
"""

from __future__ import annotations

import json
import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from myrm_agent_harness.toolkits.file_parsers import LegacyFormatParser, PptxParser, get_parser, is_supported


class TestPptxParserRegistry:
    """Test parser registration and discovery."""

    def test_pptx_is_supported(self) -> None:
        assert is_supported("slides.pptx") is True

    def test_ppt_is_supported(self) -> None:
        assert is_supported("old.ppt") is True

    def test_get_parser_returns_pptx_parser(self) -> None:
        parser = get_parser("report.pptx")
        assert isinstance(parser, PptxParser)

    def test_supported_extensions(self) -> None:
        parser = PptxParser()
        assert ".pptx" in parser.supported_extensions

    def test_ppt_uses_legacy_parser(self) -> None:
        parser = get_parser("old.ppt")
        assert isinstance(parser, LegacyFormatParser)


class TestPptxParserBasic:
    """Test basic slide text extraction."""

    @pytest.mark.asyncio
    async def test_single_slide_text(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Hello World"
            slide.placeholders[1].text = "Body content here"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser()
            result = await parser.parse(tmp)
            assert "## Slide 1" in result
            assert "Hello World" in result
            assert "Body content here" in result
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_multiple_slides(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            for i in range(3):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = f"Slide {i + 1} Title"
                slide.placeholders[1].text = f"Content {i + 1}"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser()
            result = await parser.parse(tmp)
            assert "## Slide 1" in result
            assert "## Slide 2" in result
            assert "## Slide 3" in result
            assert "Slide 1 Title" in result
            assert "Content 3" in result
            assert "---" in result
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_empty_presentation(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser()
            result = await parser.parse(tmp)
            assert result == "(Empty presentation)"
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_file_not_found(self) -> None:
        parser = PptxParser()
        with pytest.raises(FileNotFoundError):
            await parser.parse("/nonexistent/path.pptx")

    @pytest.mark.asyncio
    async def test_default_format_is_markdown(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Default Format"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser()
            result = await parser.parse(tmp)
            assert "## Slide 1" in result
            assert "Default Format" in result
        finally:
            os.unlink(tmp)


class TestPptxParserTable:
    """Test table extraction."""

    @pytest.mark.asyncio
    async def test_table_extraction(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
            table = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
            table.cell(0, 0).text = "Name"
            table.cell(0, 1).text = "Score"
            table.cell(1, 0).text = "Alice"
            table.cell(1, 1).text = "95"
            table.cell(2, 0).text = "Bob"
            table.cell(2, 1).text = "88"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser()
            result = await parser.parse(tmp)
            assert "| Name | Score |" in result
            assert "| --- | --- |" in result
            assert "| Alice | 95 |" in result
            assert "| Bob | 88 |" in result
        finally:
            os.unlink(tmp)


class TestPptxParserNotes:
    """Test speaker notes extraction."""

    @pytest.mark.asyncio
    async def test_notes_extraction(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Slide with Notes"
            slide.placeholders[1].text = "Main content"
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Remember to mention the Q2 numbers"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser()
            result = await parser.parse(tmp)
            assert "Remember to mention the Q2 numbers" in result
            assert "> **Notes:**" in result
        finally:
            os.unlink(tmp)


class TestPptxParserStructure:
    """Test structure mode JSON metadata output."""

    @pytest.mark.asyncio
    async def test_returns_valid_json(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Test Title"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser(output_format="structure")
            result = await parser.parse(tmp)
            data = json.loads(result)
            assert "slides" in data
            assert "slide_count" in data
            assert "slide_size_cm" in data
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_slide_metadata(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            for i in range(3):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = f"Slide {i + 1}"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser(output_format="structure")
            result = await parser.parse(tmp)
            data = json.loads(result)
            assert data["slide_count"] == 3
            for slide in data["slides"]:
                assert "slide_number" in slide
                assert "layout" in slide
                assert "shape_count" in slide
                assert "shapes" in slide
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_shape_ids_and_types(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Q4 Results"
            slide.placeholders[1].text = "Revenue grew 20%"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser(output_format="structure")
            result = await parser.parse(tmp)
            data = json.loads(result)
            shapes = data["slides"][0]["shapes"]
            assert len(shapes) >= 2

            types = {s["type"] for s in shapes}
            assert "title" in types

            for shape in shapes:
                assert "shape_id" in shape
                assert isinstance(shape["shape_id"], int)
                assert "name" in shape
                assert "position_cm" in shape
                pos = shape["position_cm"]
                assert "left" in pos
                assert "top" in pos
                assert "width" in pos
                assert "height" in pos
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_table_shape_in_structure(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(4), Inches(2))
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser(output_format="structure")
            result = await parser.parse(tmp)
            data = json.loads(result)
            shapes = data["slides"][0]["shapes"]
            table_shapes = [s for s in shapes if s["type"] == "table"]
            assert len(table_shapes) == 1
            assert table_shapes[0]["table_size"]["rows"] == 3
            assert table_shapes[0]["table_size"]["cols"] == 2
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_text_preview_truncated(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "A" * 300
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser(output_format="structure")
            result = await parser.parse(tmp)
            data = json.loads(result)
            title_shape = next(
                s for s in data["slides"][0]["shapes"] if s["type"] == "title"
            )
            assert len(title_shape["text_preview"]) == 200
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_notes_in_structure(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "With Notes"
            notes = slide.notes_slide
            notes.notes_text_frame.text = "Speaker notes content"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser(output_format="structure")
            result = await parser.parse(tmp)
            data = json.loads(result)
            assert "notes_preview" in data["slides"][0]
            assert "Speaker notes" in data["slides"][0]["notes_preview"]
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_empty_presentation_structure(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser(output_format="structure")
            result = await parser.parse(tmp)
            data = json.loads(result)
            assert data["slide_count"] == 0
            assert data["slides"] == []
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_structure_is_compact(self) -> None:
        """Structure mode output should be much smaller than markdown for complex presentations."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            for i in range(10):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = f"Slide {i + 1} Title"
                slide.placeholders[1].text = f"Detailed content for slide {i + 1} " * 20
            prs.save(f.name)
            tmp = f.name

        try:
            md_parser = PptxParser(output_format="markdown")
            st_parser = PptxParser(output_format="structure")
            md_result = await md_parser.parse(tmp)
            st_result = await st_parser.parse(tmp)
            assert len(st_result) < len(md_result)
        finally:
            os.unlink(tmp)


class TestPptxParserSync:
    """Test _parse_sync directly for mention.py integration."""

    def test_parse_sync_works(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Sync Test"
            prs.save(f.name)
            tmp = f.name

        try:
            parser = PptxParser()
            result = parser._parse_sync(tmp)
            assert "Sync Test" in result
        finally:
            os.unlink(tmp)
