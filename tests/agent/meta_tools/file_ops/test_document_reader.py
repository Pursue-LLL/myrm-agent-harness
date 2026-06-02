"""Tests for document_reader module

Tests is_document_path() and read_document_as_text() for .docx, .xlsx, .xls, .pptx, .ppt.
"""

from __future__ import annotations

import os
import tempfile
from unittest.mock import patch

import pytest

from myrm_agent_harness.agent.meta_tools.file_ops.utils.document_reader import (
    DOCUMENT_EXTENSIONS,
    is_document_path,
    read_document_as_text,
)


class TestIsDocumentPath:
    """Tests for is_document_path()"""

    def test_docx_detected(self) -> None:
        assert is_document_path("report.docx") is True

    def test_xlsx_detected(self) -> None:
        assert is_document_path("data.xlsx") is True

    def test_xls_detected(self) -> None:
        assert is_document_path("legacy.xls") is True

    def test_case_insensitive(self) -> None:
        assert is_document_path("REPORT.DOCX") is True
        assert is_document_path("data.XLSX") is True

    def test_pdf_not_detected(self) -> None:
        assert is_document_path("document.pdf") is False

    def test_txt_not_detected(self) -> None:
        assert is_document_path("notes.txt") is False

    def test_image_not_detected(self) -> None:
        assert is_document_path("photo.jpg") is False

    def test_no_extension(self) -> None:
        assert is_document_path("README") is False

    def test_path_with_directory(self) -> None:
        assert is_document_path("/workspace/reports/quarterly.xlsx") is True

    def test_pptx_detected(self) -> None:
        assert is_document_path("slides.pptx") is True

    def test_ppt_detected(self) -> None:
        assert is_document_path("old_slides.ppt") is True

    def test_extensions_frozen(self) -> None:
        assert isinstance(DOCUMENT_EXTENSIONS, frozenset)
        assert ".docx" in DOCUMENT_EXTENSIONS
        assert ".xlsx" in DOCUMENT_EXTENSIONS
        assert ".xls" in DOCUMENT_EXTENSIONS
        assert ".pptx" in DOCUMENT_EXTENSIONS
        assert ".ppt" in DOCUMENT_EXTENSIONS


class MockExecutor:
    """Minimal mock for CodeExecutor.read_file_bytes"""

    def __init__(self, file_data: dict[str, bytes]) -> None:
        self._data = file_data

    async def read_file_bytes(self, path: str) -> bytes:
        if path not in self._data:
            raise FileNotFoundError(f"File not found: {path}")
        return self._data[path]


class TestReadDocumentAsText:
    """Tests for read_document_as_text()"""

    @pytest.mark.asyncio
    async def test_docx_parsing(self) -> None:
        """Test .docx file is parsed to markdown via DocxParser"""
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            doc.add_heading("Test Heading", level=1)
            doc.add_paragraph("Hello World")
            doc.save(f.name)
            f.seek(0)
            docx_bytes = open(f.name, "rb").read()
            os.unlink(f.name)

        executor = MockExecutor({"contract.docx": docx_bytes})
        result = await read_document_as_text("contract.docx", executor)

        assert "[Document: contract.docx]" in result
        assert "Test Heading" in result
        assert "Hello World" in result

    @pytest.mark.asyncio
    async def test_xlsx_parsing(self) -> None:
        """Test .xlsx file is parsed to markdown table via ExcelParser"""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        assert ws is not None
        ws.append(["Name", "Age"])
        ws.append(["Alice", 30])
        ws.append(["Bob", 25])

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb.save(f.name)
            xlsx_bytes = open(f.name, "rb").read()
            os.unlink(f.name)

        executor = MockExecutor({"data.xlsx": xlsx_bytes})
        result = await read_document_as_text("data.xlsx", executor)

        assert "[Document: data.xlsx]" in result
        assert "Name" in result
        assert "Alice" in result
        assert "30" in result

    @pytest.mark.asyncio
    async def test_file_not_found(self) -> None:
        """Test FileNotFoundError is raised for missing files"""
        executor = MockExecutor({})
        with pytest.raises(FileNotFoundError):
            await read_document_as_text("missing.docx", executor)

    @pytest.mark.asyncio
    async def test_read_error_returns_message(self) -> None:
        """Test graceful handling of read errors"""

        class FailingExecutor:
            async def read_file_bytes(self, path: str) -> bytes:
                raise OSError("Disk read error")

        executor = FailingExecutor()
        result = await read_document_as_text("broken.docx", executor)
        assert "[Document: broken.docx]" in result
        assert "Failed to read" in result

    @pytest.mark.asyncio
    async def test_empty_docx(self) -> None:
        """Test empty document returns appropriate message"""
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            doc.save(f.name)
            f.seek(0)
            empty_bytes = open(f.name, "rb").read()
            os.unlink(f.name)

        executor = MockExecutor({"empty.docx": empty_bytes})
        result = await read_document_as_text("empty.docx", executor)
        assert "No extractable content" in result

    @pytest.mark.asyncio
    async def test_unsupported_extension(self) -> None:
        """Test unsupported extension returns error message"""
        executor = MockExecutor({"file.odt": b"fake"})
        result = await read_document_as_text("file.odt", executor)
        assert "Unsupported document format" in result

    @pytest.mark.asyncio
    async def test_pptx_parsing(self) -> None:
        """Test .pptx file is parsed to markdown via PptxParser"""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Quarterly Review"
            slide.placeholders[1].text = "Revenue increased by 25%"
            prs.save(f.name)
            f.seek(0)
            pptx_bytes = open(f.name, "rb").read()
            os.unlink(f.name)

        executor = MockExecutor({"slides.pptx": pptx_bytes})
        result = await read_document_as_text("slides.pptx", executor)

        assert "[Document: slides.pptx]" in result
        assert "Quarterly Review" in result
        assert "Revenue increased by 25%" in result

    @pytest.mark.asyncio
    async def test_truncation_protection(self) -> None:
        """Test large documents are truncated at _FALLBACK_MAX_CHARS"""
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            for i in range(5000):
                doc.add_paragraph(f"Paragraph {i} with some padding text to make it longer " * 3)
            doc.save(f.name)
            f.seek(0)
            large_bytes = open(f.name, "rb").read()
            os.unlink(f.name)

        executor = MockExecutor({"large.docx": large_bytes})
        result = await read_document_as_text("large.docx", executor)
        assert "truncated" in result

    @pytest.mark.asyncio
    async def test_temp_file_cleanup(self) -> None:
        """Test temp files are cleaned up after parsing"""
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            doc.add_paragraph("Cleanup test")
            doc.save(f.name)
            f.seek(0)
            cleanup_bytes = open(f.name, "rb").read()
            os.unlink(f.name)

        executor = MockExecutor({"cleanup.docx": cleanup_bytes})
        with patch("myrm_agent_harness.agent.meta_tools.file_ops.utils.document_reader.os.unlink") as mock_unlink:
            await read_document_as_text("cleanup.docx", executor)

            # Verify os.unlink was called with a .docx temp file
            assert mock_unlink.call_count == 1
            args, _ = mock_unlink.call_args
            assert args[0].endswith(".docx")
