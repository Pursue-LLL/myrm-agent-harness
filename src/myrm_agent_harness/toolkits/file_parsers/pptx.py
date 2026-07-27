"""PowerPoint document parser

Uses python-pptx for parsing PPTX files with support for:
- Slide text extraction (titles, body, text boxes)
- Table extraction with Markdown output
- Speaker notes extraction
- Structure mode: JSON metadata with shape IDs, positions, types for incremental edits

[INPUT]
- (none)

[OUTPUT]
- PptxParser: PowerPoint document parser using python-pptx

[POS]
PowerPoint document parser
"""

from __future__ import annotations

import asyncio
import json
import logging
from pathlib import Path
from typing import Literal

from myrm_agent_harness.toolkits.file_parsers.base import FileParser

logger = logging.getLogger(__name__)

PptxOutputFormat = Literal["markdown", "structure"]

_EMU_PER_CM = 360_000


def _emu_to_cm(emu: int | None) -> float | None:
    """Convert EMU (English Metric Units) to centimeters."""
    return round(emu / _EMU_PER_CM, 2) if emu is not None else None


class PptxParser(FileParser):
    """PowerPoint document parser using python-pptx

    Extracts slide text, tables, and speaker notes into Markdown format.
    Each slide becomes a section with heading and content.
    Structure mode outputs JSON metadata with shape IDs, positions, and types
    for enabling precise incremental edits.
    """

    def __init__(self, output_format: PptxOutputFormat = "markdown"):
        self._output_format: PptxOutputFormat = output_format

    async def parse(self, file_path: str) -> str:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        content = await asyncio.to_thread(self._parse_sync, file_path)

        logger.info("PowerPoint parsed: %s, format=%s, length: %d chars", path.name, self._output_format, len(content))
        return content

    def _parse_sync(self, file_path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError("python-pptx is not installed. Run: uv add python-pptx") from e

        prs = Presentation(file_path)

        if self._output_format == "structure":
            return self._build_structure(prs)

        return self._build_markdown(prs)

    def _build_markdown(self, prs: object) -> str:
        """Build Markdown content from presentation."""
        slides_output: list[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):  # type: ignore[attr-defined]
            slide_parts: list[str] = [f"## Slide {slide_idx}"]

            text_parts: list[str] = []
            table_parts: list[str] = []

            for shape in slide.shapes:
                if shape.has_table:
                    table_parts.append(self._extract_table(shape.table))
                elif shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)

            if text_parts:
                slide_parts.append("\n".join(text_parts))

            if table_parts:
                slide_parts.extend(table_parts)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> **Notes:** {notes_text}")

            if len(slide_parts) > 1:
                slides_output.append("\n\n".join(slide_parts))

        return "\n\n---\n\n".join(slides_output) if slides_output else "(Empty presentation)"

    def _build_structure(self, prs: object) -> str:
        """Return JSON structural metadata for the presentation."""
        slides_meta: list[dict[str, object]] = []

        width_cm = _emu_to_cm(prs.slide_width)  # type: ignore[attr-defined]
        height_cm = _emu_to_cm(prs.slide_height)  # type: ignore[attr-defined]

        for slide_idx, slide in enumerate(prs.slides, start=1):  # type: ignore[attr-defined]
            layout_name = slide.slide_layout.name if slide.slide_layout else None
            shapes_meta: list[dict[str, object]] = []

            for shape in slide.shapes:
                shape_info: dict[str, object] = {
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "type": self._classify_shape(shape),
                    "position_cm": {
                        "left": _emu_to_cm(shape.left),
                        "top": _emu_to_cm(shape.top),
                        "width": _emu_to_cm(shape.width),
                        "height": _emu_to_cm(shape.height),
                    },
                }

                try:
                    ph_fmt = shape.placeholder_format
                    if ph_fmt is not None:
                        shape_info["placeholder_idx"] = ph_fmt.idx
                except ValueError:
                    pass

                if shape.has_text_frame:
                    full_text = shape.text_frame.text.strip()
                    shape_info["text_preview"] = full_text[:200] if full_text else ""

                if shape.has_table:
                    tbl = shape.table
                    shape_info["table_size"] = {
                        "rows": len(tbl.rows),
                        "cols": len(tbl.columns),
                    }

                if shape.shape_type is not None and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    shape_info["has_image"] = True

                shapes_meta.append(shape_info)

            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            slide_meta: dict[str, object] = {
                "slide_number": slide_idx,
                "layout": layout_name,
                "shape_count": len(shapes_meta),
                "shapes": shapes_meta,
            }
            if notes_text:
                slide_meta["notes_preview"] = notes_text[:200]

            slides_meta.append(slide_meta)

        result: dict[str, object] = {
            "slide_count": len(slides_meta),
            "slide_size_cm": {"width": width_cm, "height": height_cm},
            "slides": slides_meta,
        }
        return json.dumps(result, ensure_ascii=False, separators=(",", ":"))

    @staticmethod
    def _classify_shape(shape: object) -> str:
        """Classify shape into a human-readable type string."""
        if hasattr(shape, "has_table") and shape.has_table:  # type: ignore[attr-defined]
            return "table"
        if hasattr(shape, "has_chart") and shape.has_chart:  # type: ignore[attr-defined]
            return "chart"
        shape_type = getattr(shape, "shape_type", None)
        if shape_type is not None and shape_type == 13:
            return "picture"
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:  # type: ignore[attr-defined]
            try:
                ph = shape.placeholder_format  # type: ignore[attr-defined]
                if ph is not None:
                    idx = ph.idx
                    if idx == 0:
                        return "title"
                    if idx == 1:
                        return "body"
                    return "placeholder"
            except (ValueError, AttributeError):
                pass
            return "text_box"
        return "other"

    @staticmethod
    def _extract_table(table: object) -> str:
        """Extract table as Markdown."""
        rows: list[list[str]] = []
        for row in table.rows:  # type: ignore[attr-defined]
            cells: list[str] = []
            for cell in row.cells:
                text = cell.text.replace("|", "\\|").replace("\n", " ").strip()
                cells.append(text)
            rows.append(cells)

        if not rows:
            return ""

        lines: list[str] = []
        headers = rows[0]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in rows[1:]:
            while len(row) < len(headers):
                row.append("")
            lines.append("| " + " | ".join(row) + " |")

        return "\n".join(lines)

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]
